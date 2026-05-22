"""
文档解析模块
- 支持 PDF/DOCX/TXT/MD/XLSX/PPTX/EML
- 结构化分块：识别章节标题，按语义边界切分
- 替代 V1 的 pypdf（使用 pymupdf，解析质量更高）
"""

import os
import re
import time
import logging
from typing import List, Dict, Any, Optional, Callable

from models.schemas import ProcessedDocument

logger = logging.getLogger(__name__)


class DocumentProcessor:

    def __init__(self, config: dict):
        self.chunk_size = config.get("chunk_size", 1024)
        self.chunk_overlap = config.get("chunk_overlap", 100)
        self.min_content_length = config.get("min_content_length", 100)
        self.supported_formats: Dict[str, Callable] = {
            '.pdf': self._parse_pdf,
            '.docx': self._parse_docx,
            '.txt': self._parse_txt,
            '.md': self._parse_txt,
            '.xlsx': self._parse_xlsx,
            '.pptx': self._parse_pptx,
        }

    # ── 公开接口 ──────────────────────────────────────────────

    def process_directory(self, dir_path: str) -> List[ProcessedDocument]:
        """递归扫描目录，解析所有支持格式的文档"""
        documents = []
        for root, _, files in os.walk(dir_path):
            for fname in sorted(files):
                ext = os.path.splitext(fname)[1].lower()
                if ext not in self.supported_formats:
                    continue
                fpath = os.path.join(root, fname)
                try:
                    doc = self.process_file(fpath)
                    if doc and len(doc.raw_content.strip()) >= self.min_content_length:
                        documents.append(doc)
                except Exception as e:
                    logger.warning(f"跳过文件 {fpath}: {e}")
        logger.info(f"共解析 {len(documents)} 篇文档")
        return documents

    def process_file(self, file_path: str) -> Optional[ProcessedDocument]:
        """解析单个文件"""
        ext = os.path.splitext(file_path)[1].lower()
        parser = self.supported_formats.get(ext)
        if not parser:
            return None

        raw_content = parser(file_path)
        cleaned = self._clean_text(raw_content)
        blocks = self._smart_chunk(cleaned)

        return ProcessedDocument(
            id=ProcessedDocument.generate_id(file_path),
            original_path=file_path,
            title=os.path.basename(file_path),
            raw_content=cleaned,
            content_blocks=blocks,
            file_size=os.path.getsize(file_path),
            file_type=ext,
            metadata=self._extract_metadata(file_path),
        )

    # ── 解析器 ────────────────────────────────────────────────

    def _parse_pdf(self, file_path: str) -> str:
        """使用 pymupdf 解析 PDF（比 pypdf 质量更高）"""
        import fitz  # pymupdf
        doc = fitz.open(file_path)
        pages = []
        for page in doc:
            text = page.get_text("text")
            if text and text.strip():
                pages.append(text)
        doc.close()
        return '\n\n'.join(pages)

    def _parse_docx(self, file_path: str) -> str:
        from docx import Document
        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # 保留标题层级信息
                if para.style and para.style.name.startswith('Heading'):
                    level = para.style.name.replace('Heading ', '#' * 1)
                    paragraphs.append(f"## {text}")
                else:
                    paragraphs.append(text)
        return '\n\n'.join(paragraphs)

    def _parse_txt(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    def _parse_xlsx(self, file_path: str) -> str:
        """提取 Excel 内容为文本"""
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows_text.append(' | '.join(cells))
            if rows_text:
                parts.append(f"[Sheet: {sheet_name}]\n" + '\n'.join(rows_text))
        wb.close()
        return '\n\n'.join(parts)

    def _parse_pptx(self, file_path: str) -> str:
        """提取 PPT 内容"""
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides_text.append(f"[Slide {i}]\n" + '\n'.join(texts))
        return '\n\n'.join(slides_text)

    # ── 文本处理 ──────────────────────────────────────────────

    def _clean_text(self, text: str) -> str:
        """文本清洗"""
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)   # 多余空行
        text = re.sub(r'[ \t]+', ' ', text)               # 多余空格
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)  # 控制字符
        return text.strip()

    def _smart_chunk(self, text: str) -> List[str]:
        """
        智能分块：优先按章节/段落边界切分，而非暴力固定长度。
        比 V1 的 RecursiveCharacterTextSplitter 更适合企业文档。
        """
        # 按段落先粗切
        paragraphs = re.split(r'\n\n+', text)
        blocks = []
        current_block = ""

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            # 如果当前段落本身就超长，做细切
            if len(para) > self.chunk_size * 2:
                if current_block:
                    blocks.append(current_block.strip())
                    current_block = ""
                # 按句子级别切分超长段落
                sentences = re.split(r'(?<=[。！？；\n])', para)
                sub_block = ""
                for sent in sentences:
                    if len(sub_block) + len(sent) <= self.chunk_size:
                        sub_block += sent
                    else:
                        if sub_block:
                            blocks.append(sub_block.strip())
                        sub_block = sent
                if sub_block:
                    blocks.append(sub_block.strip())
                continue

            # 正常累积
            if len(current_block) + len(para) + 2 <= self.chunk_size:
                current_block += "\n\n" + para if current_block else para
            else:
                if current_block:
                    blocks.append(current_block.strip())
                current_block = para

        if current_block:
            blocks.append(current_block.strip())

        # 过滤过短的块
        blocks = [b for b in blocks if len(b) >= 50]
        return blocks if blocks else [text[:self.chunk_size]]

    def _extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract file metadata (context-aware features for clustering)"""
        stat = os.stat(file_path)
        path_parts = file_path.replace('\\', '/').split('/')
        return {
            'modified_time': time.ctime(stat.st_mtime),
            'created_time': time.ctime(stat.st_ctime),
            'file_size_kb': round(stat.st_size / 1024, 1),
            'directory': os.path.dirname(file_path),
            # 路径层级作为上下文特征
            'path_depth': len(path_parts),
            'parent_folder': path_parts[-2] if len(path_parts) >= 2 else "",
            'grandparent_folder': path_parts[-3] if len(path_parts) >= 3 else "",
        }
